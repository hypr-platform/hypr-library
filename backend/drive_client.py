"""
Drive API wrapper para Biblioteca HYPR.

A Cloud Function roda como a SA biblioteca-hypr@site-hypr.iam...
A pasta foi compartilhada com essa SA como Viewer.
"""
import logging
import time
from typing import Iterator, Optional

from google.auth import default
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload
import io

log = logging.getLogger("biblioteca.drive")

# MIME types
MIME_FOLDER = "application/vnd.google-apps.folder"
MIME_SLIDES = "application/vnd.google-apps.presentation"
MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MIME_PDF = "application/pdf"

# Valid deck mime types
DECK_MIMES = {MIME_SLIDES, MIME_PPTX, MIME_PDF}

# Cutoff: só processa decks de 2025+ (criados OU modificados)
CUTOFF_DATE = "2025-01-01T00:00:00Z"


def _is_recent(file: dict) -> bool:
    """Retorna True se o arquivo foi criado OU modificado em 2025+."""
    created = file.get("createdTime", "")
    modified = file.get("modifiedTime", "")
    # ISO strings comparam lexicograficamente quando começam com YYYY-MM-DD
    return created >= CUTOFF_DATE or modified >= CUTOFF_DATE


class DriveClient:
    def __init__(self):
        credentials, _ = default(
            scopes=["https://www.googleapis.com/auth/drive.readonly"]
        )
        self._creds = credentials
        self.service = build("drive", "v3", credentials=credentials, cache_discovery=False)
        # Slides API aceita o scope drive.readonly — não precisa de scope extra,
        # só habilitar a API no projeto (ver SETUP.md ETAPA 2).
        self._slides = None

    @property
    def slides(self):
        if self._slides is None:
            self._slides = build("slides", "v1", credentials=self._creds, cache_discovery=False)
        return self._slides

    def _with_retry(self, func, max_retries=3):
        """Retry com backoff exponencial pra Drive API rate limits."""
        for attempt in range(max_retries):
            try:
                return func()
            except HttpError as e:
                if e.resp.status in (403, 429, 500, 502, 503):
                    wait = 2 ** attempt
                    log.warning(f"Drive API erro {e.resp.status}, retry em {wait}s")
                    time.sleep(wait)
                    continue
                raise
        raise RuntimeError(f"Drive API falhou após {max_retries} tentativas")

    def list_subfolders(self, parent_id: str) -> list[dict]:
        """Lista subpastas diretas de um folder (= clientes da raiz)."""
        folders = []
        page_token = None

        while True:
            def call():
                return self.service.files().list(
                    q=f"'{parent_id}' in parents and mimeType='{MIME_FOLDER}' and trashed=false",
                    fields="nextPageToken, files(id, name, modifiedTime)",
                    pageSize=200,
                    pageToken=page_token,
                    orderBy="name",
                ).execute()

            result = self._with_retry(call)
            folders.extend(result.get("files", []))
            page_token = result.get("nextPageToken")
            if not page_token:
                break

        return folders

    def list_files_in_folder(self, folder_id: str) -> list[dict]:
        """Lista todos arquivos (não-pastas) dentro de um folder, filtrando por data."""
        files = []
        page_token = None
        mime_query = " or ".join([f"mimeType='{m}'" for m in DECK_MIMES])

        # Query do Drive já filtra na fonte — economiza bandwidth e tempo
        date_filter = (
            f"(createdTime >= '{CUTOFF_DATE}' or modifiedTime >= '{CUTOFF_DATE}')"
        )

        while True:
            def call():
                return self.service.files().list(
                    q=f"'{folder_id}' in parents and ({mime_query}) and {date_filter} and trashed=false",
                    fields=(
                        "nextPageToken, files("
                        "id, name, mimeType, size, createdTime, modifiedTime, "
                        "webViewLink, thumbnailLink, iconLink, "
                        "owners(emailAddress, displayName)"
                        ")"
                    ),
                    pageSize=200,
                    pageToken=page_token,
                    orderBy="modifiedTime desc",
                ).execute()

            result = self._with_retry(call)
            files.extend(result.get("files", []))
            page_token = result.get("nextPageToken")
            if not page_token:
                break

        # Guarda extra de segurança caso a query do Drive não filtre corretamente
        return [f for f in files if _is_recent(f)]

    def iter_all_decks(self, root_folder_id: str) -> Iterator[tuple[str, dict]]:
        """
        Itera sobre TODOS os decks de TODAS as pastas-cliente.
        Yield: (client_name, file_metadata)
        """
        clients = self.list_subfolders(root_folder_id)
        log.info(f"Encontradas {len(clients)} pastas-cliente")

        for client_folder in clients:
            client_name = client_folder["name"]
            client_folder_id = client_folder["id"]
            files = self.list_files_in_folder(client_folder_id)
            log.info(f"  {client_name}: {len(files)} decks")

            for f in files:
                # Anexa client info
                f["_client"] = client_name
                f["_client_folder_id"] = client_folder_id
                yield client_name, f

    def extract_text_from_slides(self, file_id: str) -> str:
        """
        Extrai texto de um Google Slides via export pra text/plain.
        Funciona pra MIME_SLIDES. Pra .pptx/PDF, fallback diferente.
        """
        def call():
            return self.service.files().export(
                fileId=file_id,
                mimeType="text/plain",
            ).execute()

        try:
            content = self._with_retry(call)
            if isinstance(content, bytes):
                return content.decode("utf-8", errors="ignore")
            return str(content)
        except Exception as e:
            log.warning(f"Falha ao extrair texto de {file_id}: {e}")
            return ""

    def extract_text_from_pptx(self, file_id: str) -> str:
        """Download .pptx e extrai texto via python-pptx."""
        try:
            from pptx import Presentation

            buf = io.BytesIO()
            req = self.service.files().get_media(fileId=file_id)
            downloader = MediaIoBaseDownload(buf, req)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buf.seek(0)

            prs = Presentation(buf)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception as e:
            log.warning(f"Falha ao extrair pptx {file_id}: {e}")
            return ""

    # ------------------------------------------------------------
    # Extração POR SLIDE (usada pelo tagging)
    # Retorna [{"index": 1-based, "object_id": str|None, "text": str}]
    # ------------------------------------------------------------
    @staticmethod
    def _slide_text_from_elements(elements: list) -> str:
        """Percorre pageElements (inclusive grupos e tabelas) e junta textRuns."""
        parts = []

        def walk(els):
            for el in els or []:
                if "shape" in el:
                    for te in (el["shape"].get("text") or {}).get("textElements", []):
                        run = te.get("textRun")
                        if run and run.get("content"):
                            parts.append(run["content"])
                elif "table" in el:
                    for row in el["table"].get("tableRows", []):
                        for cell in row.get("tableCells", []):
                            for te in (cell.get("text") or {}).get("textElements", []):
                                run = te.get("textRun")
                                if run and run.get("content"):
                                    parts.append(run["content"])
                elif "elementGroup" in el:
                    walk(el["elementGroup"].get("children", []))

        walk(elements)
        return "".join(parts)

    def extract_slides_from_slides(self, file_id: str) -> list[dict]:
        """Google Slides → um item por slide, com objectId (pra deep-link #slide=id.X)."""
        def call():
            return self.slides.presentations().get(
                presentationId=file_id,
                fields="slides(objectId,pageElements)",
            ).execute()

        try:
            pres = self._with_retry(call)
        except Exception as e:
            log.warning(f"[Slides] Falha em presentations.get {file_id}: {e}")
            return []

        out = []
        for i, slide in enumerate(pres.get("slides", []), start=1):
            out.append({
                "index": i,
                "object_id": slide.get("objectId"),
                "text": self._slide_text_from_elements(slide.get("pageElements", [])),
            })
        return out

    def extract_slides_from_pptx(self, file_id: str) -> list[dict]:
        try:
            from pptx import Presentation

            buf = io.BytesIO()
            req = self.service.files().get_media(fileId=file_id)
            downloader = MediaIoBaseDownload(buf, req)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buf.seek(0)

            out = []
            for i, slide in enumerate(Presentation(buf).slides, start=1):
                texts = [sh.text for sh in slide.shapes if hasattr(sh, "text") and sh.text]
                out.append({"index": i, "object_id": None, "text": "\n".join(texts)})
            return out
        except Exception as e:
            log.warning(f"[Slides] Falha ao extrair pptx por slide {file_id}: {e}")
            return []

    def extract_slides(self, file_id: str, mime_type: str, fallback_text: str = "") -> list[dict]:
        """
        Roteador por slide. Pra PDF (e qualquer falha) usa o texto plano:
        o export text/plain separa páginas com form feed (\\f).
        """
        slides: list[dict] = []
        if mime_type == MIME_SLIDES:
            slides = self.extract_slides_from_slides(file_id)
        elif mime_type == MIME_PPTX:
            slides = self.extract_slides_from_pptx(file_id)

        if not slides and fallback_text:
            pages = [p for p in fallback_text.split("\f") if p.strip()] or [fallback_text]
            slides = [{"index": i, "object_id": None, "text": p} for i, p in enumerate(pages, 1)]
        return slides

    def extract_text(self, file_id: str, mime_type: str) -> str:
        """Roteador de extração baseado no mime type."""
        if mime_type == MIME_SLIDES:
            return self.extract_text_from_slides(file_id)
        elif mime_type == MIME_PPTX:
            return self.extract_text_from_pptx(file_id)
        elif mime_type == MIME_PDF:
            # PDFs já vêm com export plain text via Drive
            return self.extract_text_from_slides(file_id)
        return ""

    def get_thumbnail_url(self, file: dict) -> Optional[str]:
        """
        Retorna URL da thumbnail. Drive API entrega `thumbnailLink`,
        mas o link expira em ~1h. Pra prod, baixe e suba pro Cloud Storage,
        ou use Slides API pra gerar URL do slide 1.
        """
        return file.get("thumbnailLink")

    def fetch_thumbnail_bytes(self, file_id: str, size: int = 600) -> Optional[tuple[bytes, str]]:
        """
        Busca a thumbnail FRESCA de um arquivo via Drive API.
        Retorna (bytes, content_type) ou None se não houver.

        Diferente do thumbnailLink salvo (que expira em ~1h), este método
        pega o link no momento da chamada e baixa a imagem na hora.
        O parâmetro `size` controla a largura (Drive aceita =sNNN no fim da URL).
        """
        import urllib.request

        def call():
            return self.service.files().get(
                fileId=file_id,
                fields="thumbnailLink, mimeType",
            ).execute()

        try:
            meta = self._with_retry(call)
        except HttpError as e:
            log.warning(f"[Thumbnail] Erro ao buscar metadata de {file_id}: {e}")
            return None

        thumb_link = meta.get("thumbnailLink")
        if not thumb_link:
            return None

        # Drive thumbnailLink termina com '=s220' (tamanho). Troca pelo tamanho desejado.
        if "=s" in thumb_link:
            thumb_link = thumb_link.rsplit("=s", 1)[0] + f"=s{size}"

        # Baixa a imagem usando as credenciais da service account
        try:
            from google.auth.transport.requests import Request as AuthRequest
            # Garante que o token está fresco
            if not self._creds.valid:
                self._creds.refresh(AuthRequest())
            headers = {"Authorization": f"Bearer {self._creds.token}"}
            req = urllib.request.Request(thumb_link, headers=headers)
            with urllib.request.urlopen(req, timeout=10) as resp:
                content_type = resp.headers.get("Content-Type", "image/jpeg")
                data = resp.read()
                return (data, content_type)
        except Exception as e:
            log.warning(f"[Thumbnail] Falha ao baixar thumbnail de {file_id}: {e}")
            return None

    def get_preview_embed_url(self, file_id: str, mime_type: str) -> str:
        """URL pra iframe embed do preview."""
        if mime_type == MIME_SLIDES:
            return f"https://docs.google.com/presentation/d/{file_id}/preview"
        return f"https://drive.google.com/file/d/{file_id}/preview"
